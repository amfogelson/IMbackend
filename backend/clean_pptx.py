#!/usr/bin/env python3
"""
Script to clean up PowerPoint files for better extraction.
Removes duplicate layouts and simplifies complex shapes.
"""

import os
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

def clean_powerpoint(input_path, output_path):
    """Clean up a PowerPoint file for better extraction."""
    print(f"Cleaning PowerPoint: {input_path}")
    print(f"Output: {output_path}")
    
    # Load the original presentation
    prs = Presentation(input_path)
    print(f"Original presentation has {len(prs.slides)} slides")
    
    # Create a new clean presentation
    clean_prs = Presentation()
    
    # Copy each slide with simplified structure
    for i, slide in enumerate(prs.slides):
        print(f"\nProcessing slide {i+1}...")
        
        # Use a simple blank layout for all slides
        slide_layout = clean_prs.slide_layouts[6]  # Blank layout
        new_slide = clean_prs.slides.add_slide(slide_layout)
        
        # Copy shapes with simplification
        shapes_copied = 0
        for shape in slide.shapes:
            try:
                # Get basic properties
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Handle different shape types
                if hasattr(shape, 'shape_type'):
                    shape_type = shape.shape_type
                    
                    if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        # Copy text boxes
                        if hasattr(shape, 'text') and shape.text:
                            textbox = new_slide.shapes.add_textbox(left, top, width, height)
                            textbox.text = shape.text
                            
                            # Copy basic text formatting
                            if hasattr(shape, 'text_frame') and hasattr(textbox, 'text_frame'):
                                if shape.text_frame.paragraphs and textbox.text_frame.paragraphs:
                                    for j, para in enumerate(shape.text_frame.paragraphs):
                                        if j < len(textbox.text_frame.paragraphs):
                                            new_para = textbox.text_frame.paragraphs[j]
                                            if hasattr(para, 'font') and hasattr(new_para, 'font'):
                                                if hasattr(para.font, 'size') and para.font.size:
                                                    new_para.font.size = para.font.size
                                                if hasattr(para.font, 'bold') and para.font.bold:
                                                    new_para.font.bold = para.font.bold
                                                if hasattr(para.font, 'italic') and para.font.italic:
                                                    new_para.font.italic = para.font.italic
                            shapes_copied += 1
                            
                    elif shape_type == MSO_SHAPE_TYPE.GROUP:
                        # For group shapes, try to copy individual elements
                        print(f"  - Skipping complex group shape (shape {shapes_copied + 1})")
                        continue
                        
                    else:
                        # Copy other shapes as simple rectangles
                        new_shape = new_slide.shapes.add_shape(
                            MSO_SHAPE_TYPE.AUTO_SHAPE, left, top, width, height
                        )
                        
                        # Copy fill color if available
                        if hasattr(shape, 'fill') and hasattr(new_shape, 'fill'):
                            if hasattr(shape.fill, 'type') and shape.fill.type:
                                new_shape.fill.type = shape.fill.type
                                if hasattr(shape.fill, 'fore_color') and hasattr(new_shape.fill, 'fore_color'):
                                    if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                                        new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                        
                        shapes_copied += 1
                        
            except Exception as e:
                print(f"  - Error copying shape {shapes_copied + 1}: {e}")
                continue
        
        print(f"  ✅ Copied {shapes_copied} shapes from slide {i+1}")
    
    # Save the cleaned presentation
    clean_prs.save(output_path)
    print(f"\n✅ Cleaned presentation saved to: {output_path}")
    
    return output_path

def create_simple_test_presentations():
    """Create simple test presentations for each slide."""
    base_dir = os.path.dirname(__file__)
    infographics_dir = os.path.join(base_dir, "../infographics")
    master_path = os.path.join(infographics_dir, "infographics_master.pptx")
    
    if not os.path.exists(master_path):
        print(f"❌ Master PowerPoint not found: {master_path}")
        return
    
    # Create cleaned version
    cleaned_path = os.path.join(infographics_dir, "infographics_master_cleaned.pptx")
    clean_powerpoint(master_path, cleaned_path)
    
    # Create individual slide files
    prs = Presentation(master_path)
    
    for i, slide in enumerate(prs.slides):
        # Create a new presentation with just this slide
        single_prs = Presentation()
        slide_layout = single_prs.slide_layouts[6]  # Blank layout
        new_slide = single_prs.slides.add_slide(slide_layout)
        
        # Copy shapes from the original slide
        for shape in slide.shapes:
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                if hasattr(shape, 'shape_type'):
                    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        if hasattr(shape, 'text') and shape.text:
                            textbox = new_slide.shapes.add_textbox(left, top, width, height)
                            textbox.text = shape.text
                    else:
                        new_shape = new_slide.shapes.add_shape(
                            MSO_SHAPE_TYPE.AUTO_SHAPE, left, top, width, height
                        )
                        
                        # Copy fill color
                        if hasattr(shape, 'fill') and hasattr(new_shape, 'fill'):
                            if hasattr(shape.fill, 'type') and shape.fill.type:
                                new_shape.fill.type = shape.fill.type
                                if hasattr(shape.fill, 'fore_color') and hasattr(new_shape.fill, 'fore_color'):
                                    if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                                        new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
            except Exception as e:
                print(f"Error copying shape in slide {i+1}: {e}")
                continue
        
        # Save individual slide
        slide_path = os.path.join(infographics_dir, f"Picture{i+1}_simple.pptx")
        single_prs.save(slide_path)
        print(f"✅ Created simple slide file: {slide_path}")

if __name__ == "__main__":
    print("🧹 PowerPoint Cleanup Tool")
    print("=" * 50)
    
    # Create simple test presentations
    create_simple_test_presentations()
    
    print("\n🎯 Next Steps:")
    print("1. Check the created files in the infographics folder")
    print("2. Test the 'infographics_master_cleaned.pptx' file")
    print("3. Test the individual 'Picture*_simple.pptx' files")
    print("4. If they work better, update the backend to use these files") 