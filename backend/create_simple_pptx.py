#!/usr/bin/env python3
"""
Create simple PowerPoint files from scratch for each infographic.
This avoids the corruption issues from copying complex shapes.
"""

import os
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

def create_simple_powerpoint_for_slide(slide_number, output_path):
    """Create a simple PowerPoint with text boxes for each slide."""
    print(f"Creating simple PowerPoint for slide {slide_number}: {output_path}")
    
    # Create a new presentation
    prs = Presentation()
    
    # Use blank layout
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add a title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title.text = f"Infographic {slide_number}"
    
    # Add content based on slide number
    if slide_number == 1:
        # Slide 1: Simple layout with colored text boxes
        text1 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
        text1.text = "Blue Section"
        text1.fill.solid()
        text1.fill.fore_color.rgb = RGBColor(0, 171, 246)  # Blue
        
        text2 = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(2), Inches(1))
        text2.text = "Grey Section"
        text2.fill.solid()
        text2.fill.fore_color.rgb = RGBColor(40, 40, 40)  # Dark grey
        
    elif slide_number == 2:
        # Slide 2: Grid layout with colored text boxes
        colors = [
            RGBColor(0, 171, 246),    # Blue
            RGBColor(40, 40, 40),     # Dark grey
            RGBColor(214, 225, 1),    # Yellow
            RGBColor(255, 147, 29),   # Orange
            RGBColor(0, 0, 0),        # Black
            RGBColor(254, 32, 1)      # Red
        ]
        
        for i in range(6):
            row = i // 3
            col = i % 3
            x = Inches(1 + col * 2.5)
            y = Inches(2 + row * 1.5)
            
            text = slide.shapes.add_textbox(x, y, Inches(2), Inches(1))
            text.text = f"Box {i+1}"
            text.fill.solid()
            text.fill.fore_color.rgb = colors[i]
            
    elif slide_number == 3:
        # Slide 3: Mixed layout with text boxes
        # Main text box
        main_text = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(3))
        main_text.text = "Main Content Area"
        main_text.fill.solid()
        main_text.fill.fore_color.rgb = RGBColor(0, 171, 246)  # Blue
        
        # Inner text boxes
        inner1 = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(1), Inches(1))
        inner1.text = "A"
        inner1.fill.solid()
        inner1.fill.fore_color.rgb = RGBColor(40, 40, 40)  # Dark grey
        
        inner2 = slide.shapes.add_textbox(Inches(4), Inches(2.5), Inches(1), Inches(1))
        inner2.text = "B"
        inner2.fill.solid()
        inner2.fill.fore_color.rgb = RGBColor(214, 225, 1)  # Yellow
        
        # Add description
        desc_text = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(6), Inches(0.5))
        desc_text.text = "Complex Layout Example"
    
    # Save the presentation
    prs.save(output_path)
    print(f"✅ Created simple PowerPoint: {output_path}")

def create_all_simple_presentations():
    """Create simple PowerPoint files for all slides."""
    base_dir = os.path.dirname(__file__)
    infographics_dir = os.path.join(base_dir, "../infographics")
    
    # Create directory if it doesn't exist
    os.makedirs(infographics_dir, exist_ok=True)
    
    print("🎨 Creating Simple PowerPoint Files")
    print("=" * 50)
    
    # Create PowerPoint files for each slide
    for slide_number in range(1, 4):  # Slides 1, 2, 3
        output_path = os.path.join(infographics_dir, f"Picture{slide_number}_simple.pptx")
        create_simple_powerpoint_for_slide(slide_number, output_path)
    
    print("\n🎯 Summary:")
    print("✅ Picture1_simple.pptx - Basic blue/grey layout")
    print("✅ Picture2_simple.pptx - Grid layout with company colors")
    print("✅ Picture3_simple.pptx - Mixed layout example")
    print("\n📝 These files should download without corruption!")

if __name__ == "__main__":
    create_all_simple_presentations() 