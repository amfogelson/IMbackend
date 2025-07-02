#!/usr/bin/env python3
"""
Test script to analyze PowerPoint files and identify potential corruption sources.
"""

import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import xml.etree.ElementTree as ET
from copy import deepcopy

def analyze_pptx_structure(pptx_path):
    """Analyze the structure of a PowerPoint file to identify potential issues."""
    print(f"Analyzing: {pptx_path}")
    print("=" * 50)
    
    if not os.path.exists(pptx_path):
        print(f"❌ File not found: {pptx_path}")
        return
    
    try:
        prs = Presentation(pptx_path)
        print(f"✅ Successfully loaded PowerPoint with {len(prs.slides)} slides")
        
        # Analyze each slide
        for i, slide in enumerate(prs.slides):
            print(f"\n📄 Slide {i+1}:")
            print(f"   - Layout: {slide.slide_layout.name}")
            print(f"   - Shapes: {len(slide.shapes)}")
            
            # Analyze shapes
            shape_types = {}
            complex_shapes = []
            
            for j, shape in enumerate(slide.shapes):
                shape_type = shape.shape_type
                shape_types[shape_type] = shape_types.get(shape_type, 0) + 1
                
                # Check for potentially problematic shapes
                if shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.CHART]:
                    complex_shapes.append({
                        'index': j,
                        'type': shape_type,
                        'name': getattr(shape, 'name', 'Unknown')
                    })
                
                # Check for text with complex formatting
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.runs:
                            for run in paragraph.runs:
                                if hasattr(run, 'font') and run.font:
                                    if run.font.size and run.font.size.pt > 20:
                                        complex_shapes.append({
                                            'index': j,
                                            'type': 'Large Text',
                                            'name': f"Text: {run.text[:30]}..."
                                        })
            
            print(f"   - Shape types: {shape_types}")
            if complex_shapes:
                print(f"   - ⚠️  Complex shapes found:")
                for shape in complex_shapes:
                    print(f"     • {shape['type']}: {shape['name']}")
            else:
                print(f"   - ✅ No complex shapes detected")
        
        return prs
        
    except Exception as e:
        print(f"❌ Error loading PowerPoint: {e}")
        return None

def test_slide_extraction(pptx_path, slide_number):
    """Test extracting a specific slide to see if it works."""
    print(f"\n🧪 Testing slide extraction for slide {slide_number}")
    print("=" * 50)
    
    try:
        prs = Presentation(pptx_path)
        
        if slide_number > len(prs.slides):
            print(f"❌ Slide {slide_number} doesn't exist. Only {len(prs.slides)} slides available.")
            return False
        
        # Get the original slide
        original_slide = prs.slides[slide_number - 1]
        print(f"✅ Original slide loaded: {len(original_slide.shapes)} shapes")
        
        # Create a new presentation
        new_prs = Presentation()
        
        # Try to add a slide with the same layout
        try:
            slide_layout = original_slide.slide_layout
            new_slide = new_prs.slides.add_slide(slide_layout)
            print(f"✅ New slide created with layout: {slide_layout.name}")
        except Exception as e:
            print(f"❌ Failed to create slide with layout: {e}")
            return False
        
        # Try to copy shapes
        copied_shapes = 0
        failed_shapes = 0
        
        for i, shape in enumerate(original_slide.shapes):
            try:
                # Get basic properties
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Try to copy based on shape type
                if hasattr(shape, 'text') and shape.text:
                    # Text box
                    textbox = new_slide.shapes.add_textbox(left, top, width, height)
                    textbox.text = shape.text
                    copied_shapes += 1
                    print(f"   ✅ Copied text shape {i+1}")
                    
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Picture - this is often problematic
                    print(f"   ⚠️  Picture shape {i+1} - may cause corruption")
                    failed_shapes += 1
                    
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    # Chart - this is often problematic
                    print(f"   ⚠️  Chart shape {i+1} - may cause corruption")
                    failed_shapes += 1
                    
                elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    # Media - this is often problematic
                    print(f"   ⚠️  Media shape {i+1} - may cause corruption")
                    failed_shapes += 1
                    
                else:
                    # Try to copy as a generic shape
                    try:
                        new_shape = new_slide.shapes.add_shape(
                            shape.shape_type, left, top, width, height
                        )
                        copied_shapes += 1
                        print(f"   ✅ Copied generic shape {i+1}")
                    except:
                        print(f"   ❌ Failed to copy shape {i+1}")
                        failed_shapes += 1
                        
            except Exception as e:
                print(f"   ❌ Error copying shape {i+1}: {e}")
                failed_shapes += 1
        
        print(f"\n📊 Copy Results:")
        print(f"   - Successfully copied: {copied_shapes}")
        print(f"   - Failed to copy: {failed_shapes}")
        print(f"   - Total shapes: {len(original_slide.shapes)}")
        
        # Try to save the new presentation
        test_output = f"test_slide_{slide_number}.pptx"
        try:
            new_prs.save(test_output)
            print(f"✅ Test file saved: {test_output}")
            return True
        except Exception as e:
            print(f"❌ Failed to save test file: {e}")
            return False
            
    except Exception as e:
        print(f"❌ Error during slide extraction test: {e}")
        return False

def test_xml_copy_method(pptx_path, slide_number):
    """Test the XML copy method that we're using in the main code."""
    print(f"\n🔧 Testing XML copy method for slide {slide_number}")
    print("=" * 50)
    
    try:
        prs = Presentation(pptx_path)
        
        if slide_number > len(prs.slides):
            print(f"❌ Slide {slide_number} doesn't exist.")
            return False
        
        # Get the original slide
        original_slide = prs.slides[slide_number - 1]
        original_slide_xml = original_slide._element
        
        # Create a new presentation
        new_prs = Presentation()
        new_slide = new_prs.slides.add_slide(new_prs.slide_layouts[0])  # Use blank layout
        
        # Count elements before copying
        original_elements = len(list(original_slide_xml.spTree))
        print(f"📊 Original slide has {original_elements} elements")
        
        # Try to copy elements
        copied_elements = 0
        for shape_element in original_slide_xml.spTree:
            if shape_element.tag.endswith('sp') or shape_element.tag.endswith('pic') or shape_element.tag.endswith('graphicFrame'):
                try:
                    new_shape_element = deepcopy(shape_element)
                    new_slide._element.spTree.append(new_shape_element)
                    copied_elements += 1
                except Exception as e:
                    print(f"   ❌ Failed to copy element: {e}")
        
        print(f"📊 XML Copy Results:")
        print(f"   - Successfully copied: {copied_elements} elements")
        print(f"   - Failed to copy: {original_elements - copied_elements} elements")
        
        # Try to save
        test_output = f"test_xml_slide_{slide_number}.pptx"
        try:
            new_prs.save(test_output)
            print(f"✅ XML test file saved: {test_output}")
            return True
        except Exception as e:
            print(f"❌ Failed to save XML test file: {e}")
            return False
            
    except Exception as e:
        print(f"❌ Error during XML copy test: {e}")
        return False

def main():
    """Main test function."""
    # Path to your infographics master PowerPoint
    pptx_path = "../infographics/infographics_master.pptx"
    
    print("🔍 PowerPoint Corruption Analysis Tool")
    print("=" * 50)
    
    # Analyze the structure
    prs = analyze_pptx_structure(pptx_path)
    
    if not prs:
        print("\n❌ Cannot proceed with tests - PowerPoint file could not be loaded.")
        return
    
    # Test each slide
    for slide_num in range(1, min(4, len(prs.slides) + 1)):  # Test first 3 slides
        print(f"\n{'='*60}")
        print(f"🧪 COMPREHENSIVE TEST FOR SLIDE {slide_num}")
        print(f"{'='*60}")
        
        # Test regular extraction
        regular_success = test_slide_extraction(pptx_path, slide_num)
        
        # Test XML copy method
        xml_success = test_xml_copy_method(pptx_path, slide_num)
        
        print(f"\n📋 Summary for Slide {slide_num}:")
        print(f"   - Regular extraction: {'✅ Success' if regular_success else '❌ Failed'}")
        print(f"   - XML copy method: {'✅ Success' if xml_success else '❌ Failed'}")
    
    print(f"\n{'='*60}")
    print("📝 RECOMMENDATIONS:")
    print("=" * 60)
    print("1. Check the test files created (test_slide_*.pptx)")
    print("2. Open them in PowerPoint to see if they're corrupted")
    print("3. If XML method works better, use that in production")
    print("4. If both fail, the master PowerPoint may have complex elements")
    print("5. Consider simplifying the master PowerPoint or using a different approach")

if __name__ == "__main__":
    main() 